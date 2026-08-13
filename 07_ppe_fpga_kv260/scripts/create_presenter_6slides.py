#!/usr/bin/env python3
"""Create a six-slide presenter deck using the team's PPT theme."""

from copy import deepcopy
from pathlib import Path
import sys

sys.path.insert(0, "/tmp/ppe_ppt_deps")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont


SOURCE = Path("/home/sogang/다운로드/프로젝트_B팀_각자 내용 추가하세..pptx")
OUTPUT = Path("/home/sogang/ppe_fpga_ho/artifacts/PPE_B팀_성기호_최종_6장_터미널포함.pptx")
CAPTURE_DIR = Path("/home/sogang/ppe_fpga_ho/artifacts/terminal_capture")
RENDER_DIR = Path("/home/sogang/ppe_fpga_ho/artifacts/ppt_assets")

NAVY = RGBColor(14, 31, 53)
BLUE = RGBColor(35, 108, 180)
CYAN = RGBColor(33, 170, 190)
GREEN = RGBColor(50, 156, 105)
ORANGE = RGBColor(235, 139, 54)
RED = RGBColor(205, 73, 73)
LIGHT = RGBColor(242, 246, 250)
MID = RGBColor(208, 219, 231)
DARK = RGBColor(39, 48, 59)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 113, 124)


def terminal_png(source, output, title_text, max_lines=29):
    """Render a real transcript excerpt as a terminal-style PNG."""
    lines = Path(source).read_text(encoding="utf-8").splitlines()[:max_lines]
    width, height = 1500, 860
    img = Image.new("RGB", (width, height), (15, 20, 27))
    draw = ImageDraw.Draw(img)
    mono_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    font = ImageFont.truetype(mono_path, 25)
    title_font = ImageFont.truetype(mono_path, 27)
    draw.rounded_rectangle((0, 0, width - 1, height - 1), 24,
                           fill=(15, 20, 27), outline=(61, 73, 88), width=3)
    draw.rectangle((0, 0, width, 62), fill=(35, 42, 52))
    for i, color in enumerate(((255, 95, 87), (255, 189, 46), (40, 201, 64))):
        draw.ellipse((24 + i * 38, 20, 44 + i * 38, 40), fill=color)
    draw.text((150, 15), title_text, font=title_font, fill=(225, 232, 240))
    y = 82
    for line in lines:
        color = (190, 232, 174) if ("Vivado%" in line or "xsct%" in line or "$ " in line) else (214, 222, 231)
        if any(k in line for k in ("Successfully", "Complete!", "LOADED: True", "ERROR_COUNT=0", "0 Errors")):
            color = (88, 214, 141)
        draw.text((28, y), line[:96], font=font, fill=color)
        y += 26
        if y > height - 30:
            break
    output.parent.mkdir(parents=True, exist_ok=True)
    img.save(output)
    return output


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def remove_all_slides(prs):
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def add_text(slide, x, y, w, h, text, size=20, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Noto Sans CJK KR", margin=0.05,
             valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=MID, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    return shape


def title(slide, number, heading, subtitle=None):
    add_text(slide, 0.45, 0.18, 0.55, 0.45, f"{number:02d}", 15, True, WHITE,
             PP_ALIGN.CENTER)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45),
                                   Inches(0.18), Inches(0.55), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    # Bring text badge forward by recreating the number.
    add_text(slide, 0.45, 0.18, 0.55, 0.45, f"{number:02d}", 15, True, WHITE,
             PP_ALIGN.CENTER)
    add_text(slide, 1.15, 0.12, 11.55, 0.55, heading, 27, True, NAVY)
    if subtitle:
        add_text(slide, 1.17, 0.64, 11.2, 0.3, subtitle, 11, False, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.95),
                                  Inches(12.4), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = MID; line.line.fill.background()


def footer(slide, text="PPE Monitoring AI Hardware Accelerator | KV260"):
    add_text(slide, 0.48, 7.18, 10.5, 0.2, text, 8, False, GRAY)
    add_text(slide, 11.9, 7.15, 0.9, 0.22, "B TEAM", 8, True, BLUE,
             PP_ALIGN.RIGHT)


def placeholder(slide, x, y, w, h, heading, instructions):
    shape = rect(slide, x, y, w, h, LIGHT, BLUE, True)
    shape.line.dash_style = 2
    add_text(slide, x + 0.2, y + 0.22, w - 0.4, 0.35, "사진 / 스크린샷 삽입", 13,
             True, BLUE, PP_ALIGN.CENTER)
    add_text(slide, x + 0.25, y + h * 0.39, w - 0.5, 0.45, heading, 18, True,
             NAVY, PP_ALIGN.CENTER)
    add_text(slide, x + 0.35, y + h * 0.61, w - 0.7, h * 0.25, instructions,
             10, False, GRAY, PP_ALIGN.CENTER)


def chip(slide, x, y, w, text, color):
    s = rect(slide, x, y, w, 0.45, color, color, True)
    s.line.fill.background()
    add_text(slide, x + 0.05, y, w - 0.1, 0.45, text, 11, True, WHITE,
             PP_ALIGN.CENTER)


def arrow(slide, x, y, w=0.45):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                               Inches(w), Inches(0.34))
    s.fill.solid(); s.fill.fore_color.rgb = MID; s.line.fill.background()


def notes(slide, script):
    tf = slide.notes_slide.notes_text_frame
    tf.text = script


def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    title(s, 1, "나의 담당 역할 — Codex와 model.6 FPGA 통합",
          "팀 공동 YOLOv8n 프로젝트에서 C2f 하드웨어 통합·디버깅·보드 구동을 담당")
    stages = [("구조·설계도\nAI 분석", BLUE), ("HLS IP\n인터페이스", CYAN),
              ("Vivado BD\n통합", GREEN), ("병목·오류\n디버깅", ORANGE),
              ("Bitstream\n생성", RED), ("KV260\nBring-up", NAVY)]
    x = 0.55
    for i, (txt, c) in enumerate(stages):
        rect(s, x, 1.55, 1.65, 1.05, WHITE, c)
        add_text(s, x + 0.08, 1.67, 1.49, 0.8, txt, 15, True, c, PP_ALIGN.CENTER)
        if i != len(stages)-1:
            arrow(s, x + 1.71, 1.90, 0.42)
        x += 2.08
    add_picture(s, RENDER_DIR / "model6_c2f_design.png", 0.7, 3.05, 5.55, 2.75)
    rect(s, 6.7, 3.05, 5.9, 2.75, LIGHT, MID)
    add_text(s, 7.05, 3.28, 5.2, 0.42, "나의 핵심 업무", 19, True, NAVY)
    bullets = ["model.6 C2f 구조·특징맵 shape·설계 Flow 확인",
               "HLS IP 포트 및 AXI/DMA/Clock/Reset 연결",
               "Vivado 혼잡·Timing·DRC 로그 기반 반복 수정",
               "Vitis/XSCT 빌드·다운로드와 KV260 환경 검증",
               "Codex로 이미지·코드·로그·문서를 교차 분석"]
    for i,b in enumerate(bullets):
        add_text(s, 7.08, 3.84+i*0.38, 5.0, 0.31, "• "+b, 11, False, DARK)
    rect(s,0.7,6.02,11.9,0.78,NAVY,NAVY)
    add_text(s,0.95,6.08,2.2,0.28,"AI 협업 워크플로",12,True,WHITE)
    add_text(s,3.0,6.05,9.25,0.58,
             "설계 이미지 + 팀 GitHub + 로컬 파일  →  Codex 분석  →  터미널 실행  →  로그 재검증·문서화",
             11,True,WHITE,PP_ALIGN.CENTER)
    footer(s)
    notes(s, "이 프로젝트는 저를 포함한 다섯 명의 팀원이 공동으로 수행했으며 GitHub main의 전체 코드를 제가 혼자 개발한 것은 아닙니다. 제 담당은 조원들이 만든 HLS IP를 이해하고 model.6 C2f 블록으로 Vivado에 통합하는 일이었습니다. 처음에는 Codex에 프로젝트 규칙과 환경을 설정하고, 모델 구조와 설계 Flow 이미지를 보여 주어 코드의 shape와 데이터 흐름을 함께 확인했습니다. 이후 실제 GitHub와 로컬 파일, component.xml, Vivado와 Vitis 로그를 읽혀 AXI, DMA, Clock, Reset 연결과 병목 원인을 분석했습니다. 제가 Vivado 퍼센트, XSCT, KV260 SSH 터미널에서 명령을 실행하고 결과를 다시 전달하면 Codex가 다음 수정과 검증 절차를 제안하는 반복 방식으로 진행했습니다. 최종적으로 Bitstream과 KV260의 웹캠 장치 및 Overlay 로딩까지 직접 확인했습니다.")


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    title(s, 2, "Vivado Block Design — CPU와 FPGA를 하나로 연결",
          "PS는 YOLO 전·후단, PL은 model.6, DMA는 특징맵 전송을 담당")
    placeholder(s, 0.45, 1.2, 8.2, 5.55, "Vivado 최종 Block Design 전체",
                "PS·AXI DMA·model.6 IP·Clock/Reset이\n모두 보이게 Zoom to Fit로 캡처")
    rect(s, 8.95, 1.2, 3.9, 1.2, LIGHT, BLUE)
    add_text(s, 9.15, 1.35, 3.5, 0.32, "PS (ARM CPU + DDR)", 15, True, BLUE)
    add_text(s, 9.15, 1.76, 3.45, 0.38, "model.0~5 / cv2+7~22 실행", 11)
    rect(s, 8.95, 2.62, 3.9, 1.2, LIGHT, ORANGE)
    add_text(s, 9.15, 2.78, 3.5, 0.32, "AXI DMA", 15, True, ORANGE)
    add_text(s, 9.15, 3.18, 3.45, 0.38, "DDR ↔ FPGA 특징맵 전송", 11)
    rect(s, 8.95, 4.04, 3.9, 1.2, LIGHT, GREEN)
    add_text(s, 9.15, 4.20, 3.5, 0.32, "PL (FPGA Logic)", 15, True, GREEN)
    add_text(s, 9.15, 4.60, 3.45, 0.38, "model.6 cv1~Concat 실행", 11)
    rect(s, 8.95, 5.48, 3.9, 1.27, WHITE, MID)
    chip(s, 9.12, 5.68, 1.05, "AXI-Lite", BLUE)
    add_text(s, 10.28, 5.68, 2.25, 0.45, "제어·스케일", 10)
    chip(s, 9.12, 6.18, 1.05, "Stream", GREEN)
    add_text(s, 10.28, 6.18, 2.25, 0.45, "IP 간 연속 데이터", 10)
    footer(s)
    notes(s, "Vivado Block Design에서는 ARM CPU가 있는 PS와 FPGA 로직인 PL을 연결했습니다. 웹캠 영상 전체가 FPGA로 바로 들어가는 구조는 아닙니다. CPU가 YOLO model.0부터 5까지 계산한 후 128×40×40 특징맵을 DDR에 두고, AXI DMA가 이를 FPGA로 전송합니다. FPGA 내부의 HLS IP들은 AXI-Stream으로 연결되어 model.6를 처리하고, 256×40×40 결과를 DMA로 CPU에 돌려줍니다. AXI-Lite는 연산 시작과 스케일 설정을 담당합니다.")


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    title(s, 3, "실패에서 수렴까지 — 예상 밖의 병목과 판단",
          "Simulation PASS 이후 발생한 물리 배선·Stream·Memory·DMA 문제를 증거로 분해")
    headers = [(0.45, "관찰된 실패", RED), (3.47, "숫자로 확인한 원인", ORANGE), (7.42, "조정과 결과", GREEN)]
    widths = [2.82, 3.75, 5.43]
    for (x,txt,c),w in zip(headers,widths): chip(s,x,1.25,w,txt,c)
    rows = [
        ("Routing 불법", "66,590 node overlaps", "Handshake/병렬 구조 분석\n→ routing 오류 96%↓"),
        ("자원·배선 포화", "cv2 포함 LUT ≈ 91%", "cv2 CPU 이전\n→ LUT ≈ 56%"),
        ("Stream 교착", "약 82 pixel에서 정지", "DATAFLOW + FIFO\n→ 1,600 pixel 처리"),
        ("출력값 불일치", "BRAM 주소 + CHW/HWC 충돌", "stride-4 COE + HWC 계약\n→ 팀 mismatch 0"),
        ("DMA·보드 대기", "14-bit length < 409,600B", "DMA 14→26bit / IOC 판정\n→ 개인 Overlay 검증"),
    ]
    y=1.88
    for i,row in enumerate(rows):
        fill = WHITE if i%2==0 else LIGHT
        for x,w,txt in zip([0.45,3.47,7.42],widths,row):
            rect(s,x,y,w,0.82,fill,MID,False)
            add_text(s,x+0.12,y+0.05,w-0.24,0.72,txt,10,False,DARK,PP_ALIGN.CENTER)
        y += 0.91
    rect(s,0.45,6.55,12.4,0.42,NAVY,NAVY)
    add_text(s,0.65,6.55,12.0,0.42,
             "검증 범위  |  Bitstream 성공  →  KV260 Overlay 로드  ·  DMA 실데이터 End-to-End는 후속 과제",
             10,True,WHITE,PP_ALIGN.CENTER)
    footer(s)
    notes(s, "가장 어려웠던 점은 단위 IP의 simulation 성공이 전체 FPGA 성공을 보장하지 않았다는 것입니다. 전체 통합에서는 최대 6만6천590개의 node overlap과 congestion level 6이 발생했습니다. 조합 handshake를 끊고 병렬도와 LUTRAM 집중을 분석했으며, cv2를 CPU로 이동해 하드웨어 경계를 다시 정했습니다. C Simulation에서는 보이지 않던 Stream 교착은 DATAFLOW와 FIFO로 해결했고, BRAM byte와 word 주소 차이는 입력을 0으로 만든 실험으로 4c 규칙을 찾아냈습니다. DMA 길이와 HWC 배열까지 수정했지만 제 개인 보드에서는 MMIO 전체 실행이 남았습니다. 그래서 성공과 미완료 범위를 분리해 기록했습니다.")


def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    title(s, 4, "실제 개발 기록 — Vivado에서 Vitis/XSCT까지",
          "Timing·DRC·Bitstream을 확인한 뒤 KV260에 하드웨어와 ELF를 다운로드")
    vivado_img = RENDER_DIR / "vivado_terminal.png"
    vitis_img = RENDER_DIR / "vitis_terminal.png"
    add_picture(s, vivado_img, 0.45, 1.22, 6.1, 4.95)
    add_picture(s, vitis_img, 6.78, 1.22, 6.1, 4.95)
    chip(s,0.72,6.35,2.9,"Vivado  |  WNS +0.220 ns",GREEN)
    chip(s,3.83,6.35,2.45,"DRC Error 0",BLUE)
    chip(s,7.05,6.35,2.6,"Vitis  |  ELF Build",ORANGE)
    chip(s,9.86,6.35,2.72,"XSCT  |  Board Download",NAVY)
    footer(s)
    notes(s, "이 화면은 실제 개발 로그에서 가져온 기록입니다. Vivado에서는 최종 라우팅 WNS가 플러스 0.220나노초였고 DRC 오류 0, 비트스트림 생성 완료를 확인했습니다. 다음으로 Vitis 2022.2에서 하드웨어 플랫폼과 BSP를 구성하고 dma_loopback_v2 ELF를 빌드했습니다. Vitis 내부 XSCT가 FPGA 비트스트림과 XSA를 로드하고, FSBL과 ELF를 Cortex-A53에 다운로드했습니다. 다만 최종 UART DMA PASS 문자열은 보관 로그에 없기 때문에 성공으로 과장하지 않았습니다.")


def slide5(prs):
    s=prs.slides.add_slide(prs.slide_layouts[0])
    title(s,5,"KV260 보드 검증 — 직접 확인 범위와 팀 최종 목표",
          "개인 검증과 팀 공동 파이프라인 결과를 구분하여 기록")
    kv_img = RENDER_DIR / "kv260_terminal.png"
    add_picture(s,kv_img,0.45,1.18,6.05,4.7)
    placeholder(s,6.78,1.18,6.05,4.7,"팀의 실시간 PPE 실행 화면",
                "조원이 검증한 화면이면 ‘팀 공동 결과’로 표기\n본인 직접 재현 후에는 ‘개인 검증 완료’로 변경")
    metrics=[("Webcam","Logitech C270",BLUE),
             ("Video Node","/dev/video0·1",CYAN),
             ("Overlay","LOADED: True",GREEN),
             ("PYNQ IP","model6_dma",ORANGE)]
    x=0.48
    for name,val,c in metrics:
        rect(s,x,6.12,2.98,0.72,WHITE,c)
        add_text(s,x+0.12,6.18,0.9,0.25,name,9,True,c)
        add_text(s,x+0.95,6.18,1.85,0.38,val,11,True,NAVY,PP_ALIGN.CENTER)
        x+=3.1
    footer(s)
    notes(s, "제가 현재 기록으로 직접 확인한 범위는 올바른 KV260에 SSH로 접속하고 Logitech C270 웹캠과 video0, video1 장치를 확인한 것, 그리고 PYNQ에서 model.6 Bitstream을 다운로드해 LOADED True와 IP 목록을 확인한 것까지입니다. DMA MMIO 읽기에서는 커널이 대기 상태가 되어 재시작했으므로 제 개인 실시간 PPE 탐지가 완료됐다고 발표하지 않습니다. 오른쪽의 탐지 화면은 조원이 성공시킨 팀 공동 결과를 넣을 수 있지만 반드시 팀 결과라고 표기해야 합니다. 이후 제가 같은 실행을 직접 재현하면 개인 검증 완료로 갱신할 수 있습니다.")


def slide6(prs):
    s=prs.slides.add_slide(prs.slide_layouts[0])
    title(s,6,"문제 해결 방식과 남은 과제",
          "AI로 가설을 넓히되 실제 로그·수치·보드 결과로 좁혀 간 반복형 개발")
    columns=[(0.48,3.86,"문제 해결 원칙",GREEN,["성공 기준을 단계별로 정의","한 번에 한 변수만 변경","전후 수치를 동일 조건 비교","미확인 결과는 성공에서 제외"]),
             (4.73,3.86,"Codex 활용 방식",ORANGE,["설계 이미지·코드 함께 분석","Vivado/Vitis 로그 원인 분류","다음 검증 명령과 가설 관리","GitHub·문서·PPT 근거 통일"]),
             (8.98,3.86,"결과와 다음 검증",BLUE,["WNS +0.220 / DRC 0","Bitstream·XSA·Overlay 완료","DMA Loopback/Golden 비교","개인 실시간 FPS 재현"])]
    for x,w,head,c,items in [(a,b,c,d,e) for a,b,c,d,e in columns]:
        rect(s,x,1.35,w,4.25,WHITE,c)
        chip(s,x+0.18,1.55,w-0.36,head,c)
        for i,item in enumerate(items):
            add_text(s,x+0.32,2.25+i*0.68,w-0.64,0.5,"✓ "+item,12,False,DARK)
    rect(s,0.48,5.78,12.36,0.48,BLUE,BLUE)
    add_text(s,0.72,5.80,11.9,0.42,
             "실패의 가치  |  계산 오류보다 인터페이스·메모리·배선·검증 기준이 시스템 성공을 좌우",
             11,True,WHITE,PP_ALIGN.CENTER)
    rect(s,0.48,6.35,12.36,0.58,NAVY,NAVY)
    add_text(s,0.82,6.38,11.7,0.48,
             "결론  |  실패를 숨기지 않고 수치로 좁혀 Bitstream까지 수렴 — 다음은 개인 End-to-End 검증",
             12,True,WHITE,PP_ALIGN.CENTER)
    footer(s,"PPE Monitoring AI Hardware Accelerator | Project Retrospective")
    notes(s, "저는 문제를 만날 때 성공 기준을 Validate, Synthesis, Legal Routing, Timing, Bitstream, Board I/O로 나눴습니다. Codex에는 설계 이미지, GitHub와 로컬 소스, component XML과 실제 로그를 함께 제공해 원인 후보와 다음 검증 명령을 만들게 했습니다. 하지만 실행과 판단은 실제 터미널 결과로 다시 확인했습니다. 그 결과 개인 빌드는 WNS 플러스 0.220나노초, DRC 오류 0, Bitstream과 XSA, Overlay 로딩까지 수렴했습니다. 동시에 DMA와 실시간 PPE의 개인 End-to-End 검증은 남은 과제로 명시했습니다. 이 경험을 통해 실패를 감추는 것보다 재현 가능한 수치로 범위를 좁히는 것이 더 중요한 엔지니어링 역량임을 배웠습니다.")


def main():
    if not SOURCE.exists():
        raise SystemExit(f"missing source: {SOURCE}")
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    RENDER_DIR.mkdir(parents=True, exist_ok=True)
    terminal_png(CAPTURE_DIR / "01_vivado_model6_success.txt",
                 RENDER_DIR / "vivado_terminal.png", "Vivado 2022.2 — model.6 implementation", 27)
    terminal_png(CAPTURE_DIR / "02_vitis_xsct_board_run.txt",
                 RENDER_DIR / "vitis_terminal.png", "Vitis 2022.2 / XSCT — KV260 download", 29)
    terminal_png(CAPTURE_DIR / "03_kv260_camera_overlay.txt",
                 RENDER_DIR / "kv260_terminal.png", "KV260 Ubuntu — Webcam + PYNQ Overlay", 29)
    prs=Presentation(SOURCE)
    remove_all_slides(prs)
    # Use the blank layout inherited from the team theme.
    for make in (slide1,slide2,slide3,slide4,slide5,slide6):
        make(prs)
    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
